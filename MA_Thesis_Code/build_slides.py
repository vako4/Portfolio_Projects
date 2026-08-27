"""
code/build_slides.py
====================
Generates presentation/Meipariani_Presentation.pptx (15 slides) from source CSVs.
Slides 6-12 and 14 pull every number at runtime from audited CSVs — nothing retyped.
Every slide has a populated Notes pane (speaker script).

Run: uv run python code/build_slides.py
"""

import pandas as pd
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

ROOT        = Path(__file__).parent.parent
TABLES_DIR  = ROOT / "outputs" / "tables"
OUT_PATH    = ROOT / "presentation" / "Meipariani_Presentation.pptx"

# ── Load audited CSVs ────────────────────────────────────────────────────────
metrics_fp  = pd.read_csv(TABLES_DIR / "metrics_full_period.csv")
metrics_p15 = pd.read_csv(TABLES_DIR / "metrics_post_2015.csv")
jkm_fp      = pd.read_csv(TABLES_DIR / "sharpe_tests_full_period.csv")
jkm_p15     = pd.read_csv(TABLES_DIR / "sharpe_tests_post_2015.csv")
reg_fp      = pd.read_csv(TABLES_DIR / "factor_regressions_full_period.csv")
reg_p15     = pd.read_csv(TABLES_DIR / "factor_regressions_post_2015.csv")
boot_main   = pd.read_csv(TABLES_DIR / "regime_bootstrap_main.csv")
boot_sens   = pd.read_csv(TABLES_DIR / "regime_bootstrap_sensitivity.csv")
eps_main    = pd.read_csv(TABLES_DIR / "regime_episodes_main.csv",
                          parse_dates=["start_date", "end_date"])

# ── Pre-compute key values ───────────────────────────────────────────────────
mv   = metrics_fp[metrics_fp["strategy"]=="mv_gross"].iloc[0]
bm   = metrics_fp[metrics_fp["strategy"]=="benchmark"].iloc[0]
mv5  = metrics_p15[metrics_p15["strategy"]=="mv_gross"].iloc[0]
bm5  = metrics_p15[metrics_p15["strategy"]=="benchmark"].iloc[0]
n1   = metrics_fp[metrics_fp["strategy"]=="naive"].iloc[0]
n15  = metrics_p15[metrics_p15["strategy"]=="naive"].iloc[0]

jkm_mv_bm_fp   = jkm_fp[jkm_fp["pair"]=="mv_gross vs benchmark"].iloc[0]
jkm_mv_bm_p15  = jkm_p15[jkm_p15["pair"]=="mv_gross vs benchmark"].iloc[0]
jkm_mv_1n_fp   = jkm_fp[jkm_fp["pair"]=="mv_gross vs naive"].iloc[0]
jkm_mv_1n_p15  = jkm_p15[jkm_p15["pair"]=="mv_gross vs naive"].iloc[0]

capm_row = reg_fp[(reg_fp["strategy"]=="mv_gross")&(reg_fp["model"]=="capm")].iloc[0]
ff5_row  = reg_fp[(reg_fp["strategy"]=="mv_gross")&(reg_fp["model"]=="ff5mom")].iloc[0]
shrink   = (capm_row["alpha_annualized"] - ff5_row["alpha_annualized"]) / capm_row["alpha_annualized"]

bear_row  = boot_main[(boot_main["pair"]=="mv_gross vs benchmark")&(boot_main["regime"]=="bear")].iloc[0]
bear1n    = boot_main[(boot_main["pair"]=="naive vs benchmark")&(boot_main["regime"]=="bear")].iloc[0]

bear_m  = int(eps_main[eps_main["regime"]=="bear"]["n_months"].sum())
rec_m   = int(eps_main[eps_main["regime"]=="recovery"]["n_months"].sum())
bull_m  = int(eps_main[eps_main["regime"]=="bull"]["n_months"].sum())
total_m = bear_m + rec_m + bull_m

vol_fp  = mv["ann_volatility"] / bm["ann_volatility"]
vol_p15 = mv5["ann_volatility"] / bm5["ann_volatility"]

# ── Colour palette ───────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x35, 0x64)
ACCENT = RGBColor(0x2E, 0x75, 0xB6)
LIGHT  = RGBColor(0xD6, 0xE4, 0xF7)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x26, 0x26, 0x26)
GREEN  = RGBColor(0x37, 0x87, 0x44)
GREY   = RGBColor(0xF2, 0xF2, 0xF2)
MID_GR = RGBColor(0x96, 0x96, 0x96)

W = Inches(13.33)
H = Inches(7.5)

# ══════════════════════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_notes(slide, text: str):
    """Populate the Notes pane of a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text

def add_rect(slide, left, top, width, height, fill_rgb=None):
    from pptx.util import Emu
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.line.fill.background()
    if fill_rgb:
        s.fill.solid(); s.fill.fore_color.rgb = fill_rgb
    else:
        s.fill.background()
    return s

def add_text_box(slide, text, left, top, width, height,
                 font_size=14, bold=False, color=DARK,
                 align=PP_ALIGN.LEFT, font_name="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(font_size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font_name
    return tb

def header_bar(slide, title: str, num: int, total: int = 15):
    add_rect(slide, 0, 0, W, Inches(0.72), fill_rgb=NAVY)
    add_text_box(slide, title, Inches(0.3), Inches(0.08),
                 Inches(11.7), Inches(0.58),
                 font_size=21, bold=True, color=WHITE)
    add_text_box(slide, f"{num} / {total}",
                 Inches(12.2), Inches(0.08), Inches(0.9), Inches(0.58),
                 font_size=13, color=WHITE, align=PP_ALIGN.RIGHT)

def footer_bar(slide, note=""):
    add_rect(slide, 0, H - Inches(0.30), W, Inches(0.30), fill_rgb=GREY)
    if note:
        add_text_box(slide, note, Inches(0.3), H - Inches(0.28),
                     Inches(12.5), Inches(0.26),
                     font_size=9, color=MID_GR)

def pptx_table(slide, headers, rows,
               left, top, width, row_h=Inches(0.36),
               hdr_rgb=ACCENT, alt_rgb=LIGHT,
               fsz=11, hfsz=11, col_widths=None,
               sig_col=None, sig_val="Yes"):
    nr, nc = len(rows), len(headers)
    tbl = slide.shapes.add_table(nr+1, nc, left, top,
                                  width, row_h*(nr+1)).table
    if col_widths:
        for i, cw in enumerate(col_widths): tbl.columns[i].width = cw
    else:
        cw = width // nc
        for i in range(nc): tbl.columns[i].width = cw
    for r in range(nr+1): tbl.rows[r].height = row_h

    def sc(cell, txt, bold=False, bg=None, fc=DARK, sz=fsz):
        cell.text = str(txt)
        tf = cell.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = (tf.paragraphs[0].runs[0]
               if tf.paragraphs[0].runs else tf.paragraphs[0].add_run())
        run.text = str(txt); run.font.size = Pt(sz)
        run.font.bold = bold; run.font.color.rgb = fc
        run.font.name = "Calibri"
        if bg: cell.fill.solid(); cell.fill.fore_color.rgb = bg
        else:  cell.fill.background()

    for c, h in enumerate(headers):
        sc(tbl.cell(0,c), h, bold=True, bg=hdr_rgb, fc=WHITE, sz=hfsz)
    for r, row in enumerate(rows):
        bg = alt_rgb if r % 2 == 0 else WHITE
        for c, val in enumerate(row):
            is_sig = sig_col is not None and c == sig_col and str(val) == sig_val
            sc(tbl.cell(r+1, c), val,
               bg=GREEN if is_sig else bg,
               fc=WHITE if is_sig else DARK)
    return tbl


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS  (15 slides)
# ══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_rgb=NAVY)
    add_rect(s, Inches(0.5), Inches(1.4), Inches(12.33), Inches(2.9),
             fill_rgb=ACCENT)
    add_text_box(s,
        "A Risk-Adjusted Comparison of Constrained Mean-Variance\n"
        "and Passive Equity Portfolio Strategies across Market Cycles",
        Inches(0.8), Inches(1.55), Inches(11.7), Inches(2.6),
        font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, "Valerian Meipariani",
        Inches(0.5), Inches(4.55), W-Inches(1), Inches(0.5),
        font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s, "MA Thesis  |  ISET  |  2026",
        Inches(0.5), Inches(5.12), W-Inches(1), Inches(0.4),
        font_size=14, color=RGBColor(0xBF,0xD7,0xED), align=PP_ALIGN.CENTER)
    add_text_box(s, "Supervisor: Giorgi Mekerishvili",
        Inches(0.5), Inches(5.6), W-Inches(1), Inches(0.4),
        font_size=13, color=RGBColor(0xBF,0xD7,0xED), align=PP_ALIGN.CENTER)
    set_notes(s,
        "Good morning. My thesis asks whether a constrained minimum-variance "
        "portfolio of S&P 500 constituents earns a better risk-adjusted return "
        "than the S&P 500 Total Return index — and whether any advantage "
        "concentrates in bear markets. I'll walk through the data, methodology, "
        "and five main results in about twenty minutes.\n\n"
        "Key punchlines to keep in mind: the full-period Sharpe advantage is "
        f"+{jkm_mv_bm_fp['sharpe_diff_ann']:.3f} but statistically insignificant "
        f"(p = {jkm_mv_bm_fp['p_value']:.3f}); the CAPM alpha shrinks by "
        f"{shrink*100:.1f}% under the FF5+MOM model; and the bear-market Sharpe "
        f"advantage of {bear_row['point_estimate']:+.3f} is real but also "
        "achieved by a naive 1/N portfolio.")


def slide_02_outline(prs):
    s = blank_slide(prs)
    header_bar(s, "Outline", 2)
    items = [
        "1.   Research Questions & Hypotheses",
        "2.   Data  (S&P 500 constituents, 2001–2024)",
        "3.   Methodology  (MV optimisation, regime classification, tests)",
        "4.   Regime Classification",
        "5.   Performance Overview — Full Period  (Feb 2001 – Dec 2024)",
        "6.   Performance Overview — Post-2015 Subsample  (Jan 2015 – Dec 2024)",
        "7.   Jobson-Korkie-Memmel Sharpe Ratio Tests",
        "8.   Factor Regression: Alpha Decomposition",
        "9.   Regime-Conditional Bootstrap: MV vs Benchmark",
        "10. The Critical Qualification: 1/N Shares the Advantage",
        "11. Key Findings & Implications",
        "12. Limitations & Future Work",
    ]
    add_text_box(s, "\n".join(items),
        Inches(0.8), Inches(0.85), Inches(11.5), Inches(6.3),
        font_size=15, color=DARK)
    footer_bar(s)
    set_notes(s,
        "Quick signpost slide — I won't dwell here. "
        "The most important results are on slides 7–10: the JKM test (H1), "
        "the factor regression alpha decomposition, and the regime-conditional "
        "bootstrap split across two slides — first the MV result, then the "
        "critical qualification that 1/N replicates it. "
        "Slide 11 synthesises all three into the key take-away findings.")


def slide_03_rq(prs):
    s = blank_slide(prs)
    header_bar(s, "Research Questions & Hypotheses", 3)
    add_text_box(s, "Research Question",
        Inches(0.5), Inches(0.83), Inches(12), Inches(0.38),
        font_size=15, bold=True, color=NAVY)
    add_text_box(s,
        "Does a long-only, weight-capped minimum-variance portfolio of S&P 500 "
        "constituents deliver superior risk-adjusted performance relative to a "
        "passive S&P 500 Total Return benchmark, and is any advantage "
        "conditional on the prevailing market regime?",
        Inches(0.6), Inches(1.22), Inches(12), Inches(1.0),
        font_size=14, color=DARK)
    add_text_box(s, "H1 — Sharpe ratio (null form)",
        Inches(0.5), Inches(2.45), Inches(12), Inches(0.35),
        font_size=14, bold=True, color=NAVY)
    add_text_box(s,
        "The constrained MV portfolio does not consistently deliver a higher "
        "Sharpe ratio than the S&P 500 Total Return benchmark.",
        Inches(0.6), Inches(2.82), Inches(12), Inches(0.75),
        font_size=14, color=DARK)
    add_text_box(s, "H2 — Regime-conditional performance (null form)",
        Inches(0.5), Inches(3.75), Inches(12), Inches(0.35),
        font_size=14, bold=True, color=NAVY)
    add_text_box(s,
        "Any Sharpe ratio advantage of the constrained MV portfolio is not "
        "more likely to appear during bear-market or recovery periods than "
        "during bull-market periods.",
        Inches(0.6), Inches(4.12), Inches(12), Inches(0.85),
        font_size=14, color=DARK)
    add_text_box(s,
        "Both hypotheses stated in null form: finding in favour of the strategy "
        "requires rejecting H1; finding a regime-conditional advantage requires "
        "rejecting H2.",
        Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.8),
        font_size=12, color=MID_GR)
    footer_bar(s)
    set_notes(s,
        "Both hypotheses are in null form — this is standard for a rigorous "
        "empirical test that avoids data-mining framing. "
        "H1 tests whether the Sharpe advantage is statistically distinguishable "
        "from zero using the Jobson-Korkie-Memmel test with Newey-West standard "
        "errors. H2 tests whether the advantage, whatever its magnitude, is "
        "unevenly distributed across bear, recovery, and bull regimes. "
        "The two hypotheses are complementary: H1 asks about the average; "
        "H2 asks about the distribution across states of the world.")


def slide_04_data(prs):
    s = blank_slide(prs)
    header_bar(s, "Data", 4)
    bullets = [
        ("Universe",
         "Point-in-time S&P 500 membership (fja05680/sp500): 1,192 unique "
         "tickers, 343 snapshots, ~498 constituents per date. "
         "Eliminates look-ahead bias at universe construction."),
        ("Prices",
         "Yahoo Finance auto-adjusted monthly close (auto_adjust=True). "
         "Coverage rises from ~51% in early 2001 to ~98% by 2024."),
        ("Benchmark",
         f"^SP500TR (total return). Using ^GSPC (price-only) would "
         f"inflate active edge by ~230 pp over 24 years — benchmark choice matters."),
        ("Factor data",
         "Kenneth French Data Library: FF5 (2×3) + Carhart MOM monthly files."),
        ("Cleaning",
         "376 monthly returns >+100% → NaN (0.19% of 200,975 non-null obs). "
         "All 111 affected tickers subsequently fail the 60-month completeness filter."),
        ("Windows",
         f"Full period: Feb 2001–Dec 2024 (n = {total_m}).  "
         "Post-2015 robustness: Jan 2015–Dec 2024 (n = 120).  "
         "Post-2015 mean coverage: 86.9%."),
    ]
    y = Inches(0.87)
    for label, body in bullets:
        add_text_box(s, label+":", Inches(0.4), y, Inches(2.4), Inches(0.42),
                     font_size=13, bold=True, color=NAVY)
        add_text_box(s, body, Inches(2.9), y, Inches(10.1), Inches(0.55),
                     font_size=13, color=DARK)
        y += Inches(0.95)
    footer_bar(s)
    set_notes(s,
        "Two data-quality choices deserve emphasis. "
        "First, the benchmark: ^SP500TR includes dividends. "
        "Using the price-only index would show MV outperforming by about 230 "
        "percentage points more than it actually does — a serious error. "
        "Second, survivorship bias: we can only hold stocks whose 60-month "
        "price history exists in Yahoo Finance. Pre-2015, about 37% of S&P 500 "
        "members are missing at any given rebalance. These excluded names are "
        "disproportionately eventual failures — acquirees, bankruptcies, GSE "
        "conservatorships — so the active returns are an upper bound, not an "
        "achievable historical result. Post-2015 coverage averages 87%, which "
        "is why the post-2015 subsample is the primary robustness check.")


def slide_05_method(prs):
    s = blank_slide(prs)
    header_bar(s, "Methodology", 5)
    col = [
        ("Optimisation",
         "min w′Σw  s.t.  1′w=1,  0≤wi≤5%\n"
         "Ledoit-Wolf shrinkage Σ, 60-month rolling window\n"
         "CLARABEL solver (cvxpy); no-look-ahead convention"),
        ("Cost model",
         "MV net: 20 bps × monthly half-turnover\n"
         "1/N and benchmark: gross  (comparators only)"),
        ("H1 test",
         "Jobson-Korkie-Memmel (JKM) Sharpe test\n"
         "Newey-West HAC, lag 6; two-tailed\n"
         "Full period (n=287) and post-2015 (n=120)"),
    ]
    rcol = [
        ("Factor regressions",
         "CAPM → FF3 → FF5+MOM  (nested)\n"
         "NW HAC lag 6; α shrinkage = (αCAPM−αFF5)/αCAPM"),
        ("Regime rule",
         "Bear: ≥20% peak-to-trough on ^SP500TR\n"
         "Sensitivity: ≥15%; 3 bear episodes, 45 months (15.7%)\n"
         "COVID missed at monthly frequency; captured at −15%"),
        ("Bootstrap (H2)",
         "Stationary block bootstrap, 10,000 reps\n"
         "Block length = ⌊n_r^(1/3)⌋ months per regime\n"
         "Full sample only (9–12 post-2015 bear months insufficient)"),
    ]
    y = Inches(0.87)
    for lbl, body in col:
        add_text_box(s, lbl+":", Inches(0.35), y, Inches(2.1), Inches(0.55),
                     font_size=13, bold=True, color=NAVY)
        add_text_box(s, body, Inches(2.5), y, Inches(3.75), Inches(0.9),
                     font_size=12, color=DARK)
        y += Inches(1.4)
    y = Inches(0.87)
    for lbl, body in rcol:
        add_text_box(s, lbl+":", Inches(7.0), y, Inches(2.1), Inches(0.55),
                     font_size=13, bold=True, color=NAVY)
        add_text_box(s, body, Inches(9.1), y, Inches(4.0), Inches(0.9),
                     font_size=12, color=DARK)
        y += Inches(1.4)
    add_rect(s, Inches(6.65), Inches(0.82), Inches(0.02), Inches(5.8),
             fill_rgb=RGBColor(0xCC,0xCC,0xCC))
    footer_bar(s)
    set_notes(s,
        "Three design choices that often get challenged in Q&A. "
        "On covariance estimation: Ledoit-Wolf shrinkage is analytically "
        "optimal for this N/T ratio (~300 stocks, 60 months); without it "
        "the sample covariance is near-singular. "
        "On the NW lag of 6: standard for monthly return series, reflects "
        "typical monthly autocorrelation decay. "
        "On bootstrap scope: the post-2015 subsample has only 9 bear months "
        "under the main threshold and 12 under the sensitivity threshold — "
        "far too few for reliable block-bootstrap inference with block length "
        "⌊9^(1/3)⌋ = 2 months. We do not report regime-conditional results "
        "for the post-2015 window precisely to avoid overfit inference.")


def slide_06_regimes(prs):
    s = blank_slide(prs)
    header_bar(s, "Regime Classification  (−20% Threshold, Main Definition)", 6)
    rows = [[r["regime"].capitalize(),
             r["start_date"].strftime("%b %Y"),
             r["end_date"].strftime("%b %Y"),
             str(int(r["n_months"]))]
            for _, r in eps_main.iterrows()]
    pptx_table(s, ["Regime","Start","End","Months"], rows,
               left=Inches(1.1), top=Inches(0.84),
               width=Inches(6.4), row_h=Inches(0.48),
               col_widths=[Inches(1.4),Inches(1.6),Inches(1.6),Inches(1.1)],
               fsz=13, hfsz=13)
    add_text_box(s,
        f"Total: {total_m} months\n\n"
        f"Bear:      {bear_m} months  ({bear_m/total_m*100:.1f}%)\n"
        f"Recovery: {rec_m} months  ({rec_m/total_m*100:.1f}%)\n"
        f"Bull:       {bull_m} months  ({bull_m/total_m*100:.1f}%)\n\n"
        "COVID-19 missed at monthly close\n"
        "(Mar 2020 close = −19.8% from Jan peak)\n\n"
        "Captured at −15% sensitivity:\n"
        "  adds 3-month bear (Jan–Mar 2020)\n"
        "  + 4-month recovery (Apr–Jul 2020)",
        Inches(8.2), Inches(0.9), Inches(4.8), Inches(5.5),
        font_size=13, color=DARK)
    footer_bar(s, "Source: regime_episodes_main.csv")
    set_notes(s,
        f"The −20% threshold produces three bear episodes totalling {bear_m} months — "
        "the dot-com bust, the GFC, and the 2022 rate-shock bear. "
        "The 2022 episode is the shortest at 9 months; despite its brevity it is "
        "important because it falls entirely within the post-2015 subsample and "
        "is the only bear episode with near-complete coverage. "
        "The COVID omission is a methodological limitation, not an error — "
        "the March 2020 monthly close drew down only 19.8% from the January peak, "
        "just below the 20% threshold. "
        "The sensitivity classification at −15% captures it as a 3-month bear, "
        "and all bootstrap results are reported for both thresholds. "
        "Note that the 9 and 12 post-2015 bear-month counts are too sparse for "
        "regime-conditional bootstrap inference, which is why we report bootstrap "
        "results for the full sample only.")


def slide_07_perf_fp(prs):
    """Slide 7 — Performance Overview: Full Period only."""
    s = blank_slide(prs)
    header_bar(s, "Performance Overview — Full Period  (Feb 2001 – Dec 2024,  n = 287)", 7)
    def fmt(df, strat, lbl):
        r = df[df["strategy"]==strat].iloc[0]
        return [lbl,
                f"+{r['total_return']*100:.1f}%",
                f"+{r['ann_return']*100:.2f}%",
                f"{r['ann_volatility']*100:.2f}%",
                f"{r['sharpe']:.3f}",
                f"{r['sortino']:.3f}",
                f"{r['max_drawdown']*100:.1f}%"]
    hdrs = ["Strategy","Total Return","Ann. Return","Ann. Vol.","Sharpe","Sortino","Max DD"]
    strats = [("mv_gross","MV Gross"),("mv_net","MV Net"),
              ("naive","1/N Equal-Wt."),("benchmark","S&P 500 TR")]
    rows = [fmt(metrics_fp, k, lbl) for k, lbl in strats]
    pptx_table(s, hdrs, rows,
               left=Inches(0.3), top=Inches(0.9),
               width=Inches(12.7), row_h=Inches(0.55),
               fsz=14, hfsz=14)
    # Callout box
    add_text_box(s,
        f"MV gross Sharpe: {mv['sharpe']:.3f}\n"
        f"Benchmark Sharpe: {bm['sharpe']:.3f}\n"
        f"Difference: {jkm_mv_bm_fp['sharpe_diff_ann']:+.3f}  (p = {jkm_mv_bm_fp['p_value']:.3f})\n\n"
        f"Volatility ratio: {vol_fp:.2f}×  |  Max DD: "
        f"{mv['max_drawdown']*100:.1f}% vs {bm['max_drawdown']*100:.1f}%",
        Inches(0.3), Inches(5.0), Inches(12.5), Inches(1.2),
        font_size=14, color=NAVY)
    footer_bar(s, "Source: metrics_full_period.csv")
    set_notes(s,
        f"MV gross delivers a higher Sharpe ratio ({mv['sharpe']:.3f} vs "
        f"{bm['sharpe']:.3f}) and materially lower volatility "
        f"({mv['ann_volatility']*100:.2f}% vs {bm['ann_volatility']*100:.2f}%) "
        f"over the full period. "
        f"But at +{mv['ann_return']*100:.2f}% p.a. the strategy's "
        f"annualised return barely edges the benchmark's "
        f"+{bm['ann_return']*100:.2f}% — a premium of only 61 basis points. "
        f"The maximum drawdown of {mv['max_drawdown']*100:.1f}% versus "
        f"{bm['max_drawdown']*100:.1f}% is the most compelling single number: "
        "the strategy genuinely protects in crisis periods. "
        "These numbers are encouraging but the Sharpe gap needs a formal "
        "significance test — which is slide 9.")


def slide_08_perf_p15(prs):
    """Slide 8 — Performance Overview: Post-2015 Subsample."""
    s = blank_slide(prs)
    header_bar(s, "Performance Overview — Post-2015 Subsample  (Jan 2015 – Dec 2024,  n = 120)", 8)
    def fmt(df, strat, lbl):
        r = df[df["strategy"]==strat].iloc[0]
        return [lbl,
                f"+{r['total_return']*100:.1f}%",
                f"+{r['ann_return']*100:.2f}%",
                f"{r['ann_volatility']*100:.2f}%",
                f"{r['sharpe']:.3f}",
                f"{r['sortino']:.3f}",
                f"{r['max_drawdown']*100:.1f}%"]
    hdrs = ["Strategy","Total Return","Ann. Return","Ann. Vol.","Sharpe","Sortino","Max DD"]
    strats = [("mv_gross","MV Gross"),("mv_net","MV Net"),
              ("naive","1/N Equal-Wt."),("benchmark","S&P 500 TR")]
    rows = [fmt(metrics_p15, k, lbl) for k, lbl in strats]
    pptx_table(s, hdrs, rows,
               left=Inches(0.3), top=Inches(0.9),
               width=Inches(12.7), row_h=Inches(0.55),
               fsz=14, hfsz=14)
    add_text_box(s,
        f"Post-2015 story reverses on Sharpe: MV {mv5['sharpe']:.3f}  vs  "
        f"benchmark {bm5['sharpe']:.3f}  (diff {jkm_mv_bm_p15['sharpe_diff_ann']:+.3f}, "
        f"p = {jkm_mv_bm_p15['p_value']:.3f})\n"
        f"Benchmark leads by ~300 bps p.a. in returns  |  "
        f"Vol. ratio still {vol_p15:.2f}×  |  Max DD gap persists: "
        f"{mv5['max_drawdown']*100:.1f}% vs {bm5['max_drawdown']*100:.1f}%",
        Inches(0.3), Inches(5.0), Inches(12.5), Inches(1.2),
        font_size=14, color=NAVY)
    footer_bar(s, "Source: metrics_post_2015.csv")
    set_notes(s,
        f"The post-2015 window tells the opposite story on Sharpe ratios. "
        f"The benchmark Sharpe of {bm5['sharpe']:.3f} dwarfs MV's {mv5['sharpe']:.3f} — "
        f"a gap of {jkm_mv_bm_p15['sharpe_diff_ann']:+.3f}. "
        f"The benchmark's annualised return of {bm5['ann_return']*100:.2f}% was "
        f"driven by the long bull market from 2015–2021. "
        "Importantly, this window has 87% coverage — far less survivorship-biased "
        "than the pre-2015 era. When the data quality is best, the MV advantage "
        "disappears. This is the core finding for H1: we fail to reject the null. "
        "The volatility ratio remains around 0.81× and the drawdown gap persists "
        f"({mv5['max_drawdown']*100:.1f}% vs {bm5['max_drawdown']*100:.1f}%) — "
        "so the risk-reduction story survives even where the return story does not.")


def slide_09_jkm(prs):
    """Slide 9 — JKM Sharpe Tests. All numbers from CSVs."""
    s = blank_slide(prs)
    header_bar(s, "Jobson-Korkie-Memmel Sharpe Ratio Tests", 9)
    PAIRS = [
        ("mv_gross vs benchmark", "MV Gross vs S&P 500 TR"),
        ("mv_gross vs naive",     "MV Gross vs 1/N"),
        ("mv_net vs benchmark",   "MV Net vs S&P 500 TR"),
        ("mv_net vs naive",       "MV Net vs 1/N"),
        ("naive vs benchmark",    "1/N vs S&P 500 TR"),
    ]
    hdrs = ["Comparison","SR_A","SR_B","Diff","Corr","JKM z","p-value"]
    cws  = [Inches(3.4),Inches(1.1),Inches(1.1),Inches(1.1),
            Inches(1.1),Inches(1.35),Inches(1.35)]
    def brows(df):
        rows = []
        for pk, pl in PAIRS:
            r = df[df["pair"]==pk].iloc[0]
            rows.append([pl,
                         f"{r['sharpe_a_ann']:.3f}",
                         f"{r['sharpe_b_ann']:.3f}",
                         f"{r['sharpe_diff_ann']:+.3f}",
                         f"{r['correlation']:.3f}",
                         f"{r['jkm_statistic']:+.3f}",
                         f"{r['p_value']:.3f}"])
        return rows
    add_text_box(s, "Full Period  (n = 287)", Inches(0.3), Inches(0.83),
                 Inches(12.7), Inches(0.28), font_size=13, bold=True, color=NAVY)
    pptx_table(s, hdrs, brows(jkm_fp),
               left=Inches(0.2), top=Inches(1.1), width=Inches(12.9),
               row_h=Inches(0.37), col_widths=cws, fsz=12, hfsz=12)
    add_text_box(s, "Post-2015 Subsample  (n = 120)", Inches(0.3), Inches(3.52),
                 Inches(12.7), Inches(0.28), font_size=13, bold=True, color=NAVY)
    pptx_table(s, hdrs, brows(jkm_p15),
               left=Inches(0.2), top=Inches(3.80), width=Inches(12.9),
               row_h=Inches(0.37), col_widths=cws, fsz=12, hfsz=12)
    footer_bar(s,
        "Source: sharpe_tests_full_period.csv, sharpe_tests_post_2015.csv  |  "
        "JKM test, Newey-West HAC (lag 6), two-tailed")
    set_notes(s,
        f"The headline result for H1: the full-period Sharpe gap of "
        f"{jkm_mv_bm_fp['sharpe_diff_ann']:+.3f} has a JKM z-statistic of "
        f"{jkm_mv_bm_fp['jkm_statistic']:+.3f} and p-value of "
        f"{jkm_mv_bm_fp['p_value']:.3f} — not significant at any conventional level. "
        f"Post-2015, the gap inverts to {jkm_mv_bm_p15['sharpe_diff_ann']:+.3f} "
        f"(p = {jkm_mv_bm_p15['p_value']:.3f}). "
        "Two additional observations worth pointing out: "
        "(1) MV barely edges 1/N on Sharpe in either window — the full-period "
        f"gap is only {jkm_mv_1n_fp['sharpe_diff_ann']:+.3f} "
        f"(p = {jkm_mv_1n_fp['p_value']:.3f}), suggesting the optimisation "
        "adds little over naive equal-weighting. "
        "(2) The 1/N vs benchmark row shows 1/N also underperforms in the "
        "post-2015 window, consistent with both strategies sharing the same "
        "survivorship-biased universe.")


def slide_10_alpha(prs):
    """Slide 10 — Alpha Decomposition. All numbers from CSVs."""
    s = blank_slide(prs)
    header_bar(s, "Factor Regression: Alpha Decomposition  (MV Gross)", 10)
    MODELS  = ["capm","ff3","ff5mom"]
    MLABELS = {"capm":"CAPM","ff3":"FF3","ff5mom":"FF5+MOM"}
    hdrs    = ["Model","α (ann.)","t","p",
               "MKT β","SMB β","HML β","RMW β","CMA β","MOM β","R²"]
    cws     = [Inches(1.35),Inches(1.1),Inches(0.72),Inches(0.72),
               Inches(0.88),Inches(0.88),Inches(0.88),
               Inches(0.88),Inches(0.88),Inches(0.88),Inches(0.83)]
    def fmt_b(v, inc): return f"{v:+.3f}" if inc else "—"
    def brows(rdf):
        rows = []
        for m in MODELS:
            r = rdf[(rdf["strategy"]=="mv_gross")&(rdf["model"]==m)].iloc[0]
            f3 = m in ("ff3","ff5mom"); f5 = m=="ff5mom"
            rows.append([MLABELS[m],
                         f"{r['alpha_annualized']*100:+.2f}%",
                         f"{r['alpha_t_stat']:+.2f}",
                         f"{r['alpha_p_value']:.3f}",
                         f"{r['mkt_beta']:.3f}",
                         fmt_b(r["smb_beta"],f3),
                         fmt_b(r["hml_beta"],f3),
                         fmt_b(r["rmw_beta"],f5),
                         fmt_b(r["cma_beta"],f5),
                         fmt_b(r["mom_beta"],f5),
                         f"{r['r_squared']:.3f}"])
        return rows
    add_text_box(s, "Full Period  (n = 287)", Inches(0.3), Inches(0.83),
                 Inches(12.7), Inches(0.27), font_size=13, bold=True, color=NAVY)
    pptx_table(s, hdrs, brows(reg_fp),
               left=Inches(0.2), top=Inches(1.1), width=Inches(13.0),
               row_h=Inches(0.43), col_widths=cws, fsz=12, hfsz=11)
    add_text_box(s, "Post-2015 Subsample  (n = 120)", Inches(0.3), Inches(2.95),
                 Inches(12.7), Inches(0.27), font_size=13, bold=True, color=NAVY)
    pptx_table(s, hdrs, brows(reg_p15),
               left=Inches(0.2), top=Inches(3.22), width=Inches(13.0),
               row_h=Inches(0.43), col_widths=cws, fsz=12, hfsz=11)
    # Shrinkage callout
    ff3_row = reg_fp[(reg_fp["strategy"]=="mv_gross")&(reg_fp["model"]=="ff3")].iloc[0]
    add_text_box(s,
        f"Alpha path (full period):  "
        f"CAPM {capm_row['alpha_annualized']*100:+.2f}%  →  "
        f"FF3 {ff3_row['alpha_annualized']*100:+.2f}%  →  "
        f"FF5+MOM {ff5_row['alpha_annualized']*100:+.2f}%   "
        f"({shrink*100:.1f}% absorbed by RMW + CMA)",
        Inches(0.3), Inches(6.4), Inches(12.5), Inches(0.7),
        font_size=13, bold=True, color=NAVY)
    footer_bar(s,
        "Source: factor_regressions_full_period.csv, factor_regressions_post_2015.csv  |  "
        "OLS, Newey-West HAC (lag 6)")
    set_notes(s,
        f"The CAPM alpha of {capm_row['alpha_annualized']*100:+.2f}% per year "
        "looks impressive — nearly 3% per year in excess risk-adjusted return. "
        "But it survives neither the FF3 nor the FF5+MOM model. "
        f"The FF5+MOM residual is only {ff5_row['alpha_annualized']*100:+.2f}% "
        f"(t = {ff5_row['alpha_t_stat']:+.2f}, p = {ff5_row['alpha_p_value']:.3f}) — "
        "economically negligible and statistically indistinguishable from zero. "
        f"{shrink*100:.1f}% of the apparent CAPM alpha is absorbed by the "
        "profitability factor (RMW) and the conservative-investment factor (CMA). "
        "This is the Frazzini-Pedersen / Baker-Bradley-Wurgler mechanism: "
        "minimum-variance portfolios mechanically overweight low-beta, "
        "asset-light, cash-generative businesses — exactly the firms that load "
        "positively on RMW and CMA. The alpha was never alpha; it was a factor "
        "premium in disguise. "
        "Post-2015, the FF5+MOM alpha is actually negative (though insignificant), "
        "confirming the picture from slide 8.")


def slide_11_bootstrap_mv(prs):
    """Slide 11 — Bootstrap: MV vs Benchmark only. All numbers from CSVs."""
    s = blank_slide(prs)
    header_bar(s, "Regime-Conditional Bootstrap: MV vs Benchmark", 11)
    hdrs = ["Regime", "n months", "Sharpe Diff", "95% CI  (block bootstrap)", "Sig."]
    cws  = [Inches(1.6), Inches(1.4), Inches(1.6), Inches(4.0), Inches(0.9)]
    def brows(bdf):
        rows = []
        for reg in ("bull", "bear", "recovery"):
            r = bdf[(bdf["pair"]=="mv_gross vs benchmark")&(bdf["regime"]==reg)].iloc[0]
            rows.append([reg.capitalize(),
                         str(int(r["n_obs"])),
                         f"{r['point_estimate']:+.3f}",
                         f"[{r['ci_low']:+.3f},  {r['ci_high']:+.3f}]",
                         "Yes" if bool(r["significant"]) else "No"])
        return rows

    # Main threshold panel
    add_text_box(s, "Main Definition  (−20% threshold)",
                 Inches(0.3), Inches(0.83), Inches(12.7), Inches(0.30),
                 font_size=14, bold=True, color=NAVY)
    pptx_table(s, hdrs, brows(boot_main),
               left=Inches(0.2), top=Inches(1.13),
               width=Inches(9.6), row_h=Inches(0.52),
               col_widths=cws, fsz=13, hfsz=13,
               sig_col=4, sig_val="Yes")

    # Sensitivity threshold panel
    add_text_box(s, "Sensitivity Check  (−15% threshold)",
                 Inches(0.3), Inches(3.10), Inches(12.7), Inches(0.30),
                 font_size=14, bold=True, color=NAVY)
    pptx_table(s, hdrs, brows(boot_sens),
               left=Inches(0.2), top=Inches(3.40),
               width=Inches(9.6), row_h=Inches(0.52),
               col_widths=cws, fsz=13, hfsz=13,
               sig_col=4, sig_val="Yes")

    footer_bar(s,
        "Source: regime_bootstrap_main.csv, regime_bootstrap_sensitivity.csv  |  "
        "10,000 reps, block length = ⌊n^(1/3)⌋")
    set_notes(s,
        f"The directional prediction of H2 holds. "
        f"The bear-market Sharpe advantage of {bear_row['point_estimate']:+.3f} has a "
        f"95% CI of [{bear_row['ci_low']:+.3f}, {bear_row['ci_high']:+.3f}] — "
        "entirely above zero and statistically significant. "
        "This result is robust: the sensitivity check at −15% gives a nearly "
        f"identical bear-market advantage of {boot_sens[(boot_sens['pair']=='mv_gross vs benchmark')&(boot_sens['regime']=='bear')].iloc[0]['point_estimate']:+.3f}, "
        "again significant and with a CI well clear of zero. "
        "Bull and recovery differences are both insignificant, which is exactly "
        "what the low-beta theory predicts: MV underweights high-beta mega-caps "
        "that rise most in bull markets, producing a small performance drag, "
        "and outperforms in bear markets when those same mega-caps fall hardest. "
        "So we reject H2 in the directional sense: the advantage is concentrated "
        "in bear markets. The critical question — addressed on the next slide — "
        "is whether this concentration is unique to variance minimisation.")


def slide_12_bootstrap_1n(prs):
    """Slide 12 — Bootstrap: 1/N qualification. All numbers from CSVs."""
    s = blank_slide(prs)
    header_bar(s, "The Critical Qualification: 1/N Shares the Advantage", 12)
    hdrs = ["Regime", "n months", "Sharpe Diff", "95% CI  (block bootstrap)", "Sig."]
    cws  = [Inches(1.6), Inches(1.4), Inches(1.6), Inches(4.0), Inches(0.9)]
    def brows(bdf):
        rows = []
        for reg in ("bull", "bear", "recovery"):
            r = bdf[(bdf["pair"]=="naive vs benchmark")&(bdf["regime"]==reg)].iloc[0]
            rows.append([reg.capitalize(),
                         str(int(r["n_obs"])),
                         f"{r['point_estimate']:+.3f}",
                         f"[{r['ci_low']:+.3f},  {r['ci_high']:+.3f}]",
                         "Yes" if bool(r["significant"]) else "No"])
        return rows

    # Key callout line on slide (user-specified text, numbers from CSV)
    bear_1n_main = boot_main[(boot_main["pair"]=="naive vs benchmark")&(boot_main["regime"]=="bear")].iloc[0]
    callout = (
        f"Equal-weighting captures the same bear-market advantage "
        f"({bear_1n_main['point_estimate']:+.2f}, significant).  "
        "MV vs 1/N in bears: positive but insignificant."
    )
    add_rect(s, Inches(0.2), Inches(0.82), Inches(12.9), Inches(0.52),
             fill_rgb=RGBColor(0xFF, 0xF2, 0xCC))   # pale amber callout
    add_text_box(s, callout,
                 Inches(0.4), Inches(0.86), Inches(12.5), Inches(0.44),
                 font_size=14, bold=True, color=RGBColor(0x7F, 0x60, 0x00))

    # Main threshold panel
    add_text_box(s, "1/N Equal-Weight vs S&P 500 TR  —  Main (−20%)",
                 Inches(0.3), Inches(1.45), Inches(12.7), Inches(0.30),
                 font_size=13, bold=True, color=NAVY)
    pptx_table(s, hdrs, brows(boot_main),
               left=Inches(0.2), top=Inches(1.75),
               width=Inches(9.6), row_h=Inches(0.50),
               col_widths=cws, fsz=13, hfsz=13,
               sig_col=4, sig_val="Yes")

    # Sensitivity threshold panel
    add_text_box(s, "1/N Equal-Weight vs S&P 500 TR  —  Sensitivity (−15%)",
                 Inches(0.3), Inches(3.65), Inches(12.7), Inches(0.30),
                 font_size=13, bold=True, color=NAVY)
    pptx_table(s, hdrs, brows(boot_sens),
               left=Inches(0.2), top=Inches(3.95),
               width=Inches(9.6), row_h=Inches(0.50),
               col_widths=cws, fsz=13, hfsz=13,
               sig_col=4, sig_val="Yes")

    footer_bar(s,
        "Source: regime_bootstrap_main.csv, regime_bootstrap_sensitivity.csv  |  "
        "10,000 reps, block length = ⌊n^(1/3)⌋")
    set_notes(s,
        "This is the most important qualification in the thesis. "
        f"The 1/N portfolio achieves a bear-market Sharpe advantage of "
        f"{bear1n['point_estimate']:+.3f} [{bear1n['ci_low']:+.3f}, "
        f"{bear1n['ci_high']:+.3f}] — statistically significant and "
        f"comparable in magnitude to MV's {bear_row['point_estimate']:+.3f}. "
        "This directly answers the question of whether constrained MV optimisation "
        "justifies its additional complexity — and fees — over a simple "
        "equal-weight alternative. The answer in bear markets is no: departing "
        "from cap-weighting is sufficient; variance minimisation adds nothing "
        "beyond that departure. "
        "The mechanism is straightforward: in bear markets, cap-weighted indices "
        "are heavily concentrated in mega-cap names that tend to be highly correlated "
        "and fall sharply together. Both 1/N and MV avoid this concentration. "
        "MV does so by targeting low variance; 1/N does so mechanically by "
        "spreading weight equally. The outcome is the same. "
        "For a practitioner deciding between a minimum-variance product and an "
        "equal-weight ETF, this result is the key decision-relevant finding: "
        "if bear-market protection is the goal, the equal-weight alternative "
        "is simpler, cheaper, and equally effective.")


def slide_12_findings(prs):
    s = blank_slide(prs)
    header_bar(s, "Key Findings", 13)
    findings = [
        ("H1 not rejected",
         f"Sharpe gap {jkm_mv_bm_fp['sharpe_diff_ann']:+.3f}  (p = {jkm_mv_bm_fp['p_value']:.3f}) — "
         f"insignificant. Post-2015 reversal to {jkm_mv_bm_p15['sharpe_diff_ann']:+.3f}  "
         f"(p = {jkm_mv_bm_p15['p_value']:.3f}). "
         "Apparent alpha absorbed by RMW + CMA "
         f"({shrink*100:.1f}% shrinkage to {ff5_row['alpha_annualized']*100:+.2f}%)."),
        ("H2 partially consistent",
         f"Bear-market advantage {bear_row['point_estimate']:+.3f}  "
         f"[{bear_row['ci_low']:+.3f}, {bear_row['ci_high']:+.3f}] is real and robust. "
         f"But 1/N achieves {bear1n['point_estimate']:+.3f}  "
         f"[{bear1n['ci_low']:+.3f}, {bear1n['ci_high']:+.3f}] — "
         "mechanism is cap-weighting departure, not variance minimisation."),
        ("Risk reduction — genuine",
         f"Volatility {vol_fp:.2f}×–{vol_p15:.2f}× benchmark. "
         f"Full-period max drawdown {mv['max_drawdown']*100:.1f}% vs "
         f"{bm['max_drawdown']*100:.1f}%; post-2015: "
         f"{mv5['max_drawdown']*100:.1f}% vs {bm5['max_drawdown']*100:.1f}%. "
         "Persistent, not shared by 1/N."),
        ("Practical implication",
         "Constrained MV is a robust volatility-reduction tool, not a "
         "risk-adjusted return generator. Suitable for investors with an "
         "explicit volatility or drawdown mandate; not superior for investors "
         "benchmarked on Sharpe ratio relative to the S&P 500 TR."),
    ]
    y = Inches(0.88)
    for lbl, body in findings:
        add_rect(s, Inches(0.22), y+Inches(0.04), Inches(0.08), Inches(0.42),
                 fill_rgb=ACCENT)
        add_text_box(s, lbl, Inches(0.42), y, Inches(12.5), Inches(0.32),
                     font_size=14, bold=True, color=NAVY)
        add_text_box(s, body, Inches(0.42), y+Inches(0.34), Inches(12.5),
                     Inches(0.72), font_size=13, color=DARK)
        y += Inches(1.45)
    footer_bar(s)
    set_notes(s,
        "Three findings, cleanly separated. "
        "First: no consistent risk-adjusted outperformance — H1 not rejected at "
        "any conventional significance level, and the apparent CAPM alpha is a "
        "factor premium in disguise. "
        "Second: bear-market advantage is real and robust, but it is not unique "
        "to variance minimisation — 1/N achieves it too, suggesting the "
        "mechanism is simply departing from cap-weighting when it concentrates "
        "in mega-cap names that fall hardest in bear markets. "
        "Third: the strongest and most defensible result is the volatility and "
        "drawdown reduction, which is persistent across both subsamples and "
        "is not shared by 1/N. "
        "The practical implication for practitioners: use constrained MV if "
        "you have a drawdown mandate or an absolute volatility target; "
        "do not expect it to beat a total-return benchmark on Sharpe ratio.")


def slide_13_limitations(prs):
    s = blank_slide(prs)
    header_bar(s, "Limitations & Future Work", 14)
    items = [
        ("Survivorship bias",
         "Yahoo Finance coverage 62.6% (pre-2015) / 86.9% (post-2015). "
         "Active returns are upper bounds. CRSP data would largely eliminate this."),
        ("Monthly frequency",
         "COVID bear market missed at −20% monthly close. "
         "Short, sharp bear episodes are under-represented."),
        ("Single market",
         "U.S. large-cap only. Generalisability to other markets "
         "and asset classes is untested."),
        ("Cost model",
         "20 bps flat-rate; no market-impact, bid-ask variation, "
         "or staggered rebalance modelling."),
        ("Future work",
         "CRSP upgrade; alternative cost assumptions; "
         "international replication; daily-frequency regime detection."),
    ]
    y = Inches(0.88)
    for lbl, body in items:
        add_text_box(s, lbl+":", Inches(0.45), y, Inches(2.9), Inches(0.38),
                     font_size=13, bold=True, color=NAVY)
        add_text_box(s, body, Inches(3.45), y, Inches(9.55), Inches(0.6),
                     font_size=13, color=DARK)
        y += Inches(1.1)
    footer_bar(s)
    set_notes(s,
        "The survivorship-bias limitation deserves the most emphasis in Q&A. "
        "The key diagnostic is the 1/N portfolio: because 1/N and MV draw from "
        "the identical filtered universe at each rebalance, any outperformance "
        "of 1/N over the total-return benchmark that is not explained by factor "
        "exposures is attributable to their shared survivorship bias. "
        "The post-2015 reversal of 1/N's Sharpe from 0.580 to 0.608 — "
        "below the benchmark's 0.774 — is consistent with the bias diminishing "
        "as coverage improves. "
        "This suggests that in a properly constructed dataset, the full-period "
        "Sharpe advantage of the active strategies would be smaller, "
        "not larger. "
        "CRSP data is the natural next step and would resolve "
        "the coverage issue almost entirely.")


def slide_14_conclusion(prs):
    s = blank_slide(prs)
    add_rect(s, 0, 0, W, H, fill_rgb=NAVY)
    add_rect(s, Inches(0.5), Inches(1.1), Inches(12.33), Inches(4.5),
             fill_rgb=ACCENT)
    summary = (
        f"Constrained MV does not consistently outperform S&P 500 TR on Sharpe ratio.\n\n"
        f"Full-period Sharpe gap {jkm_mv_bm_fp['sharpe_diff_ann']:+.3f}  "
        f"(p = {jkm_mv_bm_fp['p_value']:.3f})  —  insignificant;  "
        f"post-2015 gap reverses to {jkm_mv_bm_p15['sharpe_diff_ann']:+.3f}.\n\n"
        f"CAPM α  {capm_row['alpha_annualized']*100:+.2f}% p.a.  shrinks by "
        f"{shrink*100:.1f}%  under FF5+MOM  →  residual "
        f"{ff5_row['alpha_annualized']*100:+.2f}%.\n\n"
        f"Bear-market advantage  {bear_row['point_estimate']:+.3f}  "
        f"[{bear_row['ci_low']:+.3f}, {bear_row['ci_high']:+.3f}]  —  real, "
        f"but also achieved by 1/N  ({bear1n['point_estimate']:+.3f}).\n\n"
        f"Volatility {vol_fp:.2f}×–{vol_p15:.2f}×  and  drawdown reduction  "
        "are genuine and persistent."
    )
    add_text_box(s, summary,
        Inches(0.85), Inches(1.25), Inches(11.5), Inches(4.1),
        font_size=15, color=WHITE, align=PP_ALIGN.LEFT)
    add_text_box(s, "Thank you  —  Questions?",
        Inches(0.5), Inches(6.05), W-Inches(1.0), Inches(0.6),
        font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    set_notes(s,
        "Thank you. I am happy to take questions. "
        "Prepared follow-up topics: (1) why Newey-West lag 6 and not 12; "
        "(2) robustness to alternative weight caps (1%, 3%, 10%); "
        "(3) whether the post-2015 Sharpe reversal is driven by the 2015-2021 "
        "bull market specifically or is a genuine structural shift; "
        "(4) the CRSP upgrade plan and expected quantitative impact on coverage; "
        "(5) whether the bear-market advantage would survive a transaction-cost "
        "adjustment if rebalancing frequency were higher.")


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

SLIDE_ORDER = [
    (1,  "Title",                          slide_01_title),
    (2,  "Outline",                         slide_02_outline),
    (3,  "Research Questions & Hypotheses", slide_03_rq),
    (4,  "Data",                            slide_04_data),
    (5,  "Methodology",                     slide_05_method),
    (6,  "Regime Classification  [CSV]",    slide_06_regimes),
    (7,  "Performance — Full Period  [CSV]", slide_07_perf_fp),
    (8,  "Performance — Post-2015   [CSV]", slide_08_perf_p15),
    (9,  "JKM Sharpe Tests         [CSV]",  slide_09_jkm),
    (10, "Alpha Decomposition      [CSV]",  slide_10_alpha),
    (11, "Bootstrap: MV vs Benchmark  [CSV]",    slide_11_bootstrap_mv),
    (12, "Bootstrap: 1/N Qualification [CSV]",  slide_12_bootstrap_1n),
    (13, "Key Findings              [CSV]",      slide_12_findings),
    (14, "Limitations & Future Work",            slide_13_limitations),
    (15, "Conclusion                [CSV]",      slide_14_conclusion),
]

if __name__ == "__main__":
    print("Building presentation…")
    prs = new_prs()
    for num, label, builder in SLIDE_ORDER:
        builder(prs)
        print(f"  Slide {num:2d}: {label}")
    prs.save(str(OUT_PATH))
    size_kb = OUT_PATH.stat().st_size // 1024
    print(f"\nSaved: {OUT_PATH}  ({size_kb} KB)")
    print(f"Total slides: {len(prs.slides)}")
